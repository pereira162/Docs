"""
RAG Document Processor - Core do Sistema de Extração e Processamento
Foco: Extrair, processar e salvar informações em formato otimizado para IA
Princípio: Local-First, Zero-Cost, Processamento Avançado
"""

import os
import json
import uuid
import hashlib
import logging
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional, Union
from datetime import datetime
import io
import requests
from urllib.parse import urlparse

# Docling imports
from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat, DocumentStream
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling.chunking import HybridChunker
from docling_core.transforms.chunker.tokenizer.huggingface import HuggingFaceTokenizer
from transformers import AutoTokenizer

# Bibliotecas para processamento local
import pdfplumber
from docx import Document as DocxDocument
import openpyxl
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract
import easyocr

# Configuração de logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class RAGProcessor:
    """
    Processador principal para extração e formatação de documentos
    Objetivo: Transformar qualquer fonte em arquivos Markdown/JSON otimizados para IA
    """
    
    def __init__(self, output_dir: str = "./rag_outputs"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        
        # Diretórios organizados
        self.markdown_dir = self.output_dir / "markdown"
        self.json_dir = self.output_dir / "json"
        self.metadata_dir = self.output_dir / "metadata"
        self.chunks_dir = self.output_dir / "chunks"
        
        for dir_path in [self.markdown_dir, self.json_dir, self.metadata_dir, self.chunks_dir]:
            dir_path.mkdir(exist_ok=True)
        
        # Configurar Docling com opções avançadas
        self.setup_docling()
        
        # EasyOCR para imagens (alternativa ao tesseract)
        try:
            self.ocr_reader = easyocr.Reader(['pt', 'en'])
            logger.info("✅ EasyOCR initialized successfully")
        except Exception as e:
            logger.warning(f"⚠️ EasyOCR initialization failed: {e}")
            self.ocr_reader = None
    
    def setup_docling(self):
        """Configurar Docling com todas as opções avançadas"""
        try:
            # Configurações avançadas do PDF
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = True
            pipeline_options.do_table_structure = True
            pipeline_options.do_picture_classification = True
            pipeline_options.generate_picture_images = True
            pipeline_options.images_scale = 2
            
            # Configurar tokenizer para chunking
            tokenizer = HuggingFaceTokenizer(
                tokenizer=AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
            )
            
            # Inicializar conversor e chunker
            self.converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
            
            self.chunker = HybridChunker(tokenizer=tokenizer)
            
            logger.info("✅ Docling configured with advanced options")
            
        except Exception as e:
            logger.error(f"❌ Docling setup failed: {e}")
            self.converter = None
            self.chunker = None
    
    def process_source(self, source: Union[str, bytes], source_type: str = "auto", filename: str = None) -> Dict:
        """
        Processar qualquer fonte (URL, arquivo local, bytes) e gerar saídas estruturadas
        """
        try:
            # Gerar ID único para o documento
            doc_id = str(uuid.uuid4())
            timestamp = datetime.now().isoformat()
            
            # Determinar tipo de fonte
            if isinstance(source, str) and (source.startswith('http') or source.startswith('https')):
                result = self._process_url(source, doc_id, timestamp)
            elif isinstance(source, bytes):
                result = self._process_bytes(source, filename or "unknown.bin", doc_id, timestamp)
            elif isinstance(source, str) and os.path.exists(source):
                result = self._process_file(source, doc_id, timestamp)
            else:
                raise ValueError(f"Unsupported source type: {type(source)}")
            
            # Salvar metadados completos
            self._save_metadata(doc_id, result)
            
            return result
            
        except Exception as e:
            logger.error(f"❌ Error processing source: {e}")
            return {
                "success": False,
                "error": str(e),
                "doc_id": doc_id if 'doc_id' in locals() else None
            }
    
    def _process_url(self, url: str, doc_id: str, timestamp: str) -> Dict:
        """Processar URL (PDF, página web, etc.)"""
        try:
            # Detectar tipo de conteúdo
            response = requests.head(url, timeout=10)
            content_type = response.headers.get('content-type', '').lower()
            
            if 'pdf' in content_type:
                # Download e processo PDF
                pdf_response = requests.get(url, timeout=30)
                return self._process_pdf_with_docling(pdf_response.content, url, doc_id, timestamp)
            
            elif 'html' in content_type or url.endswith('.html'):
                # Processar página web
                return self._process_webpage(url, doc_id, timestamp)
            
            else:
                # Tentar download genérico
                file_response = requests.get(url, timeout=30)
                filename = os.path.basename(urlparse(url).path) or "downloaded_file"
                return self._process_bytes(file_response.content, filename, doc_id, timestamp)
                
        except Exception as e:
            logger.error(f"❌ Error processing URL {url}: {e}")
            raise
    
    def _process_file(self, file_path: str, doc_id: str, timestamp: str) -> Dict:
        """Processar arquivo local"""
        path = Path(file_path)
        
        with open(path, 'rb') as f:
            content = f.read()
        
        return self._process_bytes(content, path.name, doc_id, timestamp)
    
    def _process_bytes(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar conteúdo em bytes baseado na extensão"""
        ext = Path(filename).suffix.lower()
        
        # Roteamento baseado em extensão
        if ext == '.pdf':
            return self._process_pdf_with_docling(content, filename, doc_id, timestamp)
        elif ext in ['.docx', '.doc']:
            return self._process_docx(content, filename, doc_id, timestamp)
        elif ext in ['.xlsx', '.xls']:
            return self._process_excel(content, filename, doc_id, timestamp)
        elif ext in ['.pptx', '.ppt']:
            return self._process_powerpoint(content, filename, doc_id, timestamp)
        elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            return self._process_image(content, filename, doc_id, timestamp)
        elif ext in ['.txt', '.md']:
            return self._process_text(content, filename, doc_id, timestamp)
        elif ext == '.html':
            return self._process_html(content, filename, doc_id, timestamp)
        else:
            # Fallback: tentar processar como texto
            return self._process_unknown(content, filename, doc_id, timestamp)
    
    def _process_pdf_with_docling(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar PDF usando Docling (método principal)"""
        try:
            if not self.converter:
                # Fallback para pdfplumber se Docling não estiver disponível
                return self._process_pdf_fallback(content, filename, doc_id, timestamp)
            
            # Criar stream do documento
            buf = io.BytesIO(content)
            source = DocumentStream(name=filename, stream=buf)
            
            # Converter com Docling
            result = self.converter.convert(source)
            
            if result.status.success:
                doc = result.document
                
                # Extrair texto completo em Markdown
                markdown_content = doc.export_to_markdown()
                
                # Extrair dados estruturados em JSON
                json_content = doc.export_to_dict()
                
                # Criar chunks inteligentes
                chunks = []
                chunk_metadata = []
                
                for i, chunk in enumerate(self.chunker.chunk(doc)):
                    chunk_text = chunk.text
                    chunk_context = self.chunker.contextualize(chunk=chunk)
                    
                    chunk_data = {
                        "id": f"{doc_id}_chunk_{i}",
                        "text": chunk_text,
                        "contextualized_text": chunk_context,
                        "metadata": chunk.meta.export_json_dict() if hasattr(chunk.meta, 'export_json_dict') else {}
                    }
                    
                    chunks.append(chunk_data)
                    chunk_metadata.append(chunk_data["metadata"])
                
                # Extrair metadados avançados
                metadata = self._extract_advanced_metadata(doc, filename, content)
                
                # Salvar arquivos estruturados
                files_saved = self._save_processed_files(
                    doc_id, filename, markdown_content, json_content, chunks, metadata
                )
                
                return {
                    "success": True,
                    "doc_id": doc_id,
                    "filename": filename,
                    "timestamp": timestamp,
                    "processing_method": "docling_advanced",
                    "markdown_content": markdown_content,
                    "json_content": json_content,
                    "chunks": chunks,
                    "metadata": metadata,
                    "files_saved": files_saved,
                    "stats": {
                        "total_chunks": len(chunks),
                        "total_characters": len(markdown_content),
                        "pages": metadata.get("pages", 0),
                        "tables": metadata.get("tables", 0),
                        "images": metadata.get("images", 0)
                    }
                }
            else:
                # Fallback se Docling falhar
                return self._process_pdf_fallback(content, filename, doc_id, timestamp)
                
        except Exception as e:
            logger.error(f"❌ Docling PDF processing failed: {e}")
            return self._process_pdf_fallback(content, filename, doc_id, timestamp)
    
    def _process_pdf_fallback(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Fallback para processar PDF com pdfplumber"""
        try:
            with io.BytesIO(content) as buf:
                with pdfplumber.open(buf) as pdf:
                    # Extrair texto de todas as páginas
                    text_content = []
                    tables_data = []
                    
                    for page_num, page in enumerate(pdf.pages):
                        # Texto da página
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_content.append(f"## Página {page_num + 1}\n\n{page_text}")
                        
                        # Tabelas da página
                        tables = page.extract_tables()
                        for table_num, table in enumerate(tables):
                            if table:
                                # Converter tabela para markdown
                                table_md = self._table_to_markdown(table)
                                tables_data.append({
                                    "page": page_num + 1,
                                    "table": table_num + 1,
                                    "data": table,
                                    "markdown": table_md
                                })
                    
                    # Combinar conteúdo
                    full_text = "\n\n".join(text_content)
                    
                    # Adicionar tabelas ao texto
                    if tables_data:
                        full_text += "\n\n## Tabelas Extraídas\n\n"
                        for table_info in tables_data:
                            full_text += f"### Página {table_info['page']}, Tabela {table_info['table']}\n\n"
                            full_text += table_info['markdown'] + "\n\n"
                    
                    # Criar chunks simples
                    chunks = self._create_simple_chunks(full_text, doc_id)
                    
                    # Metadados básicos
                    metadata = {
                        "filename": filename,
                        "pages": len(pdf.pages),
                        "tables": len(tables_data),
                        "processing_method": "pdfplumber_fallback",
                        "file_size": len(content)
                    }
                    
                    # Salvar arquivos
                    files_saved = self._save_processed_files(
                        doc_id, filename, full_text, {"text": full_text, "tables": tables_data}, chunks, metadata
                    )
                    
                    return {
                        "success": True,
                        "doc_id": doc_id,
                        "filename": filename,
                        "timestamp": timestamp,
                        "processing_method": "pdfplumber_fallback",
                        "markdown_content": full_text,
                        "json_content": {"text": full_text, "tables": tables_data},
                        "chunks": chunks,
                        "metadata": metadata,
                        "files_saved": files_saved
                    }
                    
        except Exception as e:
            logger.error(f"❌ PDF fallback processing failed: {e}")
            raise
    
    def _process_docx(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar documentos Word"""
        try:
            with io.BytesIO(content) as buf:
                doc = DocxDocument(buf)
                
                # Extrair texto dos parágrafos
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                
                # Extrair tabelas
                tables_data = []
                for table_num, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    
                    tables_data.append({
                        "table": table_num + 1,
                        "data": table_data,
                        "markdown": self._table_to_markdown(table_data)
                    })
                
                # Combinar conteúdo
                full_text = "\n\n".join(paragraphs)
                
                if tables_data:
                    full_text += "\n\n## Tabelas do Documento\n\n"
                    for table_info in tables_data:
                        full_text += f"### Tabela {table_info['table']}\n\n"
                        full_text += table_info['markdown'] + "\n\n"
                
                # Processar e salvar
                chunks = self._create_simple_chunks(full_text, doc_id)
                metadata = {
                    "filename": filename,
                    "paragraphs": len(paragraphs),
                    "tables": len(tables_data),
                    "processing_method": "python_docx"
                }
                
                files_saved = self._save_processed_files(
                    doc_id, filename, full_text, {"text": full_text, "tables": tables_data}, chunks, metadata
                )
                
                return {
                    "success": True,
                    "doc_id": doc_id,
                    "filename": filename,
                    "timestamp": timestamp,
                    "processing_method": "python_docx",
                    "markdown_content": full_text,
                    "chunks": chunks,
                    "metadata": metadata,
                    "files_saved": files_saved
                }
                
        except Exception as e:
            logger.error(f"❌ DOCX processing failed: {e}")
            raise
    
    def _process_excel(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar planilhas Excel"""
        try:
            with io.BytesIO(content) as buf:
                # Ler todas as planilhas
                excel_data = pd.read_excel(buf, sheet_name=None)
                
                markdown_content = f"# Planilha: {filename}\n\n"
                json_content = {"filename": filename, "sheets": {}}
                
                for sheet_name, df in excel_data.items():
                    # Converter para markdown
                    sheet_md = f"## Planilha: {sheet_name}\n\n"
                    sheet_md += df.to_markdown(index=False) + "\n\n"
                    markdown_content += sheet_md
                    
                    # Adicionar ao JSON
                    json_content["sheets"][sheet_name] = {
                        "data": df.to_dict('records'),
                        "shape": df.shape,
                        "columns": df.columns.tolist()
                    }
                
                # Processar e salvar
                chunks = self._create_simple_chunks(markdown_content, doc_id)
                metadata = {
                    "filename": filename,
                    "sheets": list(excel_data.keys()),
                    "total_rows": sum(df.shape[0] for df in excel_data.values()),
                    "processing_method": "pandas_excel"
                }
                
                files_saved = self._save_processed_files(
                    doc_id, filename, markdown_content, json_content, chunks, metadata
                )
                
                return {
                    "success": True,
                    "doc_id": doc_id,
                    "filename": filename,
                    "timestamp": timestamp,
                    "processing_method": "pandas_excel",
                    "markdown_content": markdown_content,
                    "chunks": chunks,
                    "metadata": metadata,
                    "files_saved": files_saved
                }
                
        except Exception as e:
            logger.error(f"❌ Excel processing failed: {e}")
            raise
    
    def _process_powerpoint(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar apresentações PowerPoint"""
        try:
            with io.BytesIO(content) as buf:
                prs = Presentation(buf)
                
                markdown_content = f"# Apresentação: {filename}\n\n"
                slides_data = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = f"## Slide {slide_num}\n\n"
                    slide_content = []
                    
                    # Extrair texto de todas as formas
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text)
                    
                    if slide_content:
                        slide_text += "\n".join(slide_content) + "\n\n"
                    else:
                        slide_text += "*Slide sem conteúdo de texto*\n\n"
                    
                    markdown_content += slide_text
                    slides_data.append({
                        "slide": slide_num,
                        "content": slide_content
                    })
                
                # Processar e salvar
                chunks = self._create_simple_chunks(markdown_content, doc_id)
                metadata = {
                    "filename": filename,
                    "slides": len(prs.slides),
                    "processing_method": "python_pptx"
                }
                
                json_content = {"filename": filename, "slides": slides_data}
                
                files_saved = self._save_processed_files(
                    doc_id, filename, markdown_content, json_content, chunks, metadata
                )
                
                return {
                    "success": True,
                    "doc_id": doc_id,
                    "filename": filename,
                    "timestamp": timestamp,
                    "processing_method": "python_pptx",
                    "markdown_content": markdown_content,
                    "chunks": chunks,
                    "metadata": metadata,
                    "files_saved": files_saved
                }
                
        except Exception as e:
            logger.error(f"❌ PowerPoint processing failed: {e}")
            raise
    
    def _process_image(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar imagens com OCR"""
        try:
            # Abrir imagem
            image = Image.open(io.BytesIO(content))
            
            # OCR com EasyOCR (preferido) ou Tesseract
            text_content = ""
            
            if self.ocr_reader:
                try:
                    # EasyOCR
                    results = self.ocr_reader.readtext(content)
                    text_content = "\n".join([result[1] for result in results])
                    ocr_method = "easyocr"
                except Exception as e:
                    logger.warning(f"EasyOCR failed: {e}")
                    text_content = "OCR falhou"
                    ocr_method = "failed"
            else:
                try:
                    # Tesseract fallback
                    text_content = pytesseract.image_to_string(image, lang='por+eng')
                    ocr_method = "tesseract"
                except Exception as e:
                    logger.warning(f"Tesseract failed: {e}")
                    text_content = "OCR não disponível"
                    ocr_method = "unavailable"
            
            # Criar conteúdo markdown
            markdown_content = f"# Imagem: {filename}\n\n"
            markdown_content += f"**Dimensões:** {image.size[0]}x{image.size[1]}\n"
            markdown_content += f"**Formato:** {image.format}\n\n"
            markdown_content += "## Texto Extraído (OCR)\n\n"
            markdown_content += text_content or "*Nenhum texto detectado*"
            
            # Processar e salvar
            chunks = self._create_simple_chunks(markdown_content, doc_id)
            metadata = {
                "filename": filename,
                "image_size": image.size,
                "image_format": image.format,
                "ocr_method": ocr_method,
                "processing_method": "ocr_image"
            }
            
            json_content = {
                "filename": filename,
                "image_info": {
                    "size": image.size,
                    "format": image.format,
                    "mode": image.mode
                },
                "ocr_text": text_content,
                "ocr_method": ocr_method
            }
            
            files_saved = self._save_processed_files(
                doc_id, filename, markdown_content, json_content, chunks, metadata
            )
            
            return {
                "success": True,
                "doc_id": doc_id,
                "filename": filename,
                "timestamp": timestamp,
                "processing_method": "ocr_image",
                "markdown_content": markdown_content,
                "chunks": chunks,
                "metadata": metadata,
                "files_saved": files_saved
            }
            
        except Exception as e:
            logger.error(f"❌ Image processing failed: {e}")
            raise
    
    def _process_text(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar arquivos de texto simples"""
        try:
            # Decodificar texto
            text_content = content.decode('utf-8', errors='ignore')
            
            # Se for markdown, manter formatação
            if filename.lower().endswith('.md'):
                markdown_content = text_content
            else:
                # Converter texto simples para markdown básico
                markdown_content = f"# Arquivo: {filename}\n\n```\n{text_content}\n```"
            
            # Processar e salvar
            chunks = self._create_simple_chunks(markdown_content, doc_id)
            metadata = {
                "filename": filename,
                "character_count": len(text_content),
                "line_count": text_content.count('\n'),
                "processing_method": "text_simple"
            }
            
            json_content = {
                "filename": filename,
                "content": text_content,
                "stats": {
                    "characters": len(text_content),
                    "lines": text_content.count('\n'),
                    "words": len(text_content.split())
                }
            }
            
            files_saved = self._save_processed_files(
                doc_id, filename, markdown_content, json_content, chunks, metadata
            )
            
            return {
                "success": True,
                "doc_id": doc_id,
                "filename": filename,
                "timestamp": timestamp,
                "processing_method": "text_simple",
                "markdown_content": markdown_content,
                "chunks": chunks,
                "metadata": metadata,
                "files_saved": files_saved
            }
            
        except Exception as e:
            logger.error(f"❌ Text processing failed: {e}")
            raise
    
    def _process_html(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar arquivos HTML"""
        try:
            from bs4 import BeautifulSoup
            
            # Parse HTML
            soup = BeautifulSoup(content, 'html.parser')
            
            # Extrair texto limpo
            text_content = soup.get_text(separator='\n', strip=True)
            
            # Extrair estrutura
            title = soup.find('title')
            title_text = title.get_text() if title else filename
            
            headings = []
            for tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                for heading in soup.find_all(tag):
                    headings.append(f"{tag.upper()}: {heading.get_text(strip=True)}")
            
            # Criar markdown estruturado
            markdown_content = f"# {title_text}\n\n"
            
            if headings:
                markdown_content += "## Estrutura do Documento\n\n"
                markdown_content += "\n".join(headings) + "\n\n"
            
            markdown_content += "## Conteúdo\n\n"
            markdown_content += text_content
            
            # Processar e salvar
            chunks = self._create_simple_chunks(markdown_content, doc_id)
            metadata = {
                "filename": filename,
                "title": title_text,
                "headings_count": len(headings),
                "processing_method": "beautifulsoup_html"
            }
            
            json_content = {
                "filename": filename,
                "title": title_text,
                "headings": headings,
                "content": text_content
            }
            
            files_saved = self._save_processed_files(
                doc_id, filename, markdown_content, json_content, chunks, metadata
            )
            
            return {
                "success": True,
                "doc_id": doc_id,
                "filename": filename,
                "timestamp": timestamp,
                "processing_method": "beautifulsoup_html",
                "markdown_content": markdown_content,
                "chunks": chunks,
                "metadata": metadata,
                "files_saved": files_saved
            }
            
        except Exception as e:
            logger.error(f"❌ HTML processing failed: {e}")
            raise
    
    def _process_webpage(self, url: str, doc_id: str, timestamp: str) -> Dict:
        """Processar página web"""
        try:
            response = requests.get(url, timeout=30)
            return self._process_html(response.content, f"webpage_{urlparse(url).netloc}.html", doc_id, timestamp)
        except Exception as e:
            logger.error(f"❌ Webpage processing failed: {e}")
            raise
    
    def _process_unknown(self, content: bytes, filename: str, doc_id: str, timestamp: str) -> Dict:
        """Processar arquivo de tipo desconhecido"""
        try:
            # Tentar decodificar como texto
            try:
                text_content = content.decode('utf-8', errors='ignore')
                return self._process_text(content, filename, doc_id, timestamp)
            except:
                # Se falhar, criar informações básicas
                markdown_content = f"# Arquivo Desconhecido: {filename}\n\n"
                markdown_content += f"**Tamanho:** {len(content)} bytes\n"
                markdown_content += f"**Hash MD5:** {hashlib.md5(content).hexdigest()}\n\n"
                markdown_content += "*Tipo de arquivo não suportado para extração de texto*"
                
                chunks = self._create_simple_chunks(markdown_content, doc_id)
                metadata = {
                    "filename": filename,
                    "file_size": len(content),
                    "md5_hash": hashlib.md5(content).hexdigest(),
                    "processing_method": "unknown_file"
                }
                
                json_content = {
                    "filename": filename,
                    "file_size": len(content),
                    "md5_hash": hashlib.md5(content).hexdigest(),
                    "supported": False
                }
                
                files_saved = self._save_processed_files(
                    doc_id, filename, markdown_content, json_content, chunks, metadata
                )
                
                return {
                    "success": True,
                    "doc_id": doc_id,
                    "filename": filename,
                    "timestamp": timestamp,
                    "processing_method": "unknown_file",
                    "markdown_content": markdown_content,
                    "chunks": chunks,
                    "metadata": metadata,
                    "files_saved": files_saved
                }
                
        except Exception as e:
            logger.error(f"❌ Unknown file processing failed: {e}")
            raise
    
    def _extract_advanced_metadata(self, doc, filename: str, content: bytes) -> Dict:
        """Extrair metadados avançados do documento Docling"""
        try:
            metadata = {
                "filename": filename,
                "file_size": len(content),
                "processing_method": "docling_advanced"
            }
            
            # Metadados básicos do documento
            if hasattr(doc, 'name'):
                metadata["title"] = doc.name
            
            # Contar elementos por tipo
            if hasattr(doc, 'texts'):
                texts = doc.texts
                metadata["total_elements"] = len(texts)
                
                # Categorizar elementos
                element_counts = {}
                for item in texts:
                    if hasattr(item, 'label'):
                        label = str(item.label)
                        element_counts[label] = element_counts.get(label, 0) + 1
                
                metadata["element_types"] = element_counts
                metadata["pages"] = element_counts.get('page', 0)
                metadata["tables"] = sum(1 for k in element_counts.keys() if 'table' in k.lower())
                metadata["images"] = sum(1 for k in element_counts.keys() if 'picture' in k.lower() or 'image' in k.lower())
            
            return metadata
            
        except Exception as e:
            logger.warning(f"⚠️ Could not extract advanced metadata: {e}")
            return {
                "filename": filename,
                "file_size": len(content),
                "processing_method": "docling_basic"
            }
    
    def _create_simple_chunks(self, text: str, doc_id: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict]:
        """Criar chunks simples para documentos não processados pelo Docling"""
        chunks = []
        
        if len(text) <= chunk_size:
            chunks.append({
                "id": f"{doc_id}_chunk_0",
                "text": text,
                "metadata": {"chunk_index": 0}
            })
            return chunks
        
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]
            
            # Tentar quebrar em uma fronteira de sentença
            if end < len(text):
                last_period = chunk_text.rfind('.')
                last_newline = chunk_text.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > start + chunk_size // 2:
                    chunk_text = text[start:start + break_point + 1]
                    end = start + break_point + 1
            
            chunks.append({
                "id": f"{doc_id}_chunk_{chunk_index}",
                "text": chunk_text.strip(),
                "metadata": {
                    "chunk_index": chunk_index,
                    "start_pos": start,
                    "end_pos": end
                }
            })
            
            start = end - overlap if end < len(text) else end
            chunk_index += 1
            
            if start >= len(text):
                break
        
        return chunks
    
    def _table_to_markdown(self, table_data: List[List[str]]) -> str:
        """Converter tabela para formato Markdown"""
        if not table_data or not table_data[0]:
            return "*Tabela vazia*"
        
        # Primeira linha como cabeçalho
        headers = table_data[0]
        rows = table_data[1:] if len(table_data) > 1 else []
        
        # Criar markdown
        md_table = "| " + " | ".join(str(cell) for cell in headers) + " |\n"
        md_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        
        for row in rows:
            # Preencher células vazias se necessário
            padded_row = row + [""] * (len(headers) - len(row))
            md_table += "| " + " | ".join(str(cell) for cell in padded_row[:len(headers)]) + " |\n"
        
        return md_table
    
    def _save_processed_files(self, doc_id: str, filename: str, markdown_content: str, 
                             json_content: Dict, chunks: List[Dict], metadata: Dict) -> Dict:
        """Salvar todos os arquivos processados no sistema de arquivos local"""
        files_saved = {}
        base_filename = Path(filename).stem
        
        try:
            # 1. Salvar Markdown
            md_path = self.markdown_dir / f"{base_filename}_{doc_id}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            files_saved["markdown"] = str(md_path)
            
            # 2. Salvar JSON estruturado
            json_path = self.json_dir / f"{base_filename}_{doc_id}.json"
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(json_content, f, ensure_ascii=False, indent=2)
            files_saved["json"] = str(json_path)
            
            # 3. Salvar chunks separadamente
            chunks_path = self.chunks_dir / f"{base_filename}_{doc_id}_chunks.json"
            with open(chunks_path, 'w', encoding='utf-8') as f:
                json.dump(chunks, f, ensure_ascii=False, indent=2)
            files_saved["chunks"] = str(chunks_path)
            
            # 4. Salvar metadados
            metadata_path = self.metadata_dir / f"{base_filename}_{doc_id}_metadata.json"
            with open(metadata_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, ensure_ascii=False, indent=2)
            files_saved["metadata"] = str(metadata_path)
            
            logger.info(f"✅ Files saved for document {doc_id}")
            return files_saved
            
        except Exception as e:
            logger.error(f"❌ Error saving files for {doc_id}: {e}")
            return files_saved
    
    def _save_metadata(self, doc_id: str, result: Dict):
        """Salvar metadados completos do processamento"""
        try:
            metadata_file = self.metadata_dir / f"processing_{doc_id}.json"
            with open(metadata_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2, default=str)
            logger.info(f"✅ Processing metadata saved for {doc_id}")
        except Exception as e:
            logger.error(f"❌ Error saving processing metadata: {e}")
    
    def get_processed_documents(self) -> List[Dict]:
        """Listar todos os documentos processados"""
        documents = []
        
        try:
            for metadata_file in self.metadata_dir.glob("processing_*.json"):
                try:
                    with open(metadata_file, 'r', encoding='utf-8') as f:
                        doc_data = json.load(f)
                    
                    if doc_data.get("success"):
                        documents.append({
                            "doc_id": doc_data.get("doc_id"),
                            "filename": doc_data.get("filename"),
                            "timestamp": doc_data.get("timestamp"),
                            "processing_method": doc_data.get("processing_method"),
                            "stats": doc_data.get("stats", {}),
                            "files_saved": doc_data.get("files_saved", {})
                        })
                        
                except Exception as e:
                    logger.warning(f"⚠️ Could not load metadata from {metadata_file}: {e}")
            
            return sorted(documents, key=lambda x: x.get("timestamp", ""), reverse=True)
            
        except Exception as e:
            logger.error(f"❌ Error listing processed documents: {e}")
            return []
    
    def export_document(self, doc_id: str, format: str = "markdown") -> Optional[str]:
        """Exportar documento específico"""
        try:
            base_path = None
            
            if format == "markdown":
                base_path = self.markdown_dir
                pattern = f"*_{doc_id}.md"
            elif format == "json":
                base_path = self.json_dir
                pattern = f"*_{doc_id}.json"
            elif format == "chunks":
                base_path = self.chunks_dir
                pattern = f"*_{doc_id}_chunks.json"
            else:
                raise ValueError(f"Unsupported format: {format}")
            
            # Procurar arquivo
            files = list(base_path.glob(pattern))
            if files:
                return str(files[0])
            else:
                logger.warning(f"⚠️ No file found for doc_id {doc_id} in format {format}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error exporting document {doc_id}: {e}")
            return None
    
    def export_all_documents(self, format: str = "zip") -> Optional[str]:
        """Exportar todos os documentos processados"""
        try:
            import zipfile
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            zip_path = self.output_dir / f"rag_export_{timestamp}.zip"
            
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                # Adicionar todos os arquivos organizados
                for dir_name, dir_path in [
                    ("markdown", self.markdown_dir),
                    ("json", self.json_dir),
                    ("chunks", self.chunks_dir),
                    ("metadata", self.metadata_dir)
                ]:
                    for file_path in dir_path.glob("*"):
                        if file_path.is_file():
                            arcname = f"{dir_name}/{file_path.name}"
                            zipf.write(file_path, arcname)
            
            logger.info(f"✅ All documents exported to {zip_path}")
            return str(zip_path)
            
        except Exception as e:
            logger.error(f"❌ Error exporting all documents: {e}")
            return None
